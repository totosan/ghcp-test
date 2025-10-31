import os
import re
from pptx import Presentation
from pptx.util import Pt
from datetime import datetime

class PPTExporter:
    """Handle PowerPoint export with template support"""
    
    # Placeholder mapping
    PLACEHOLDERS = {
        'PROJECT_NAME': 'project_name',
        'CLIENT': 'client_name',
        'INDUSTRY': 'industry',
        'YEAR': 'project_year',
        'CHALLENGE': 'challenge',
        'SOLUTION': 'solution',
        'OUTCOMES': 'outcomes',
        'TECHNOLOGIES': 'technologies',
        'TEAM_SIZE': 'team_size',
        'DURATION': 'duration_months',
        'PROJECT_VALUE': 'project_value',
        'TAGS': 'tags',
        'CREATED_BY': 'created_by',
    }
    
    def __init__(self, template_path):
        """Initialize with template path"""
        self.template_path = template_path
    
    def export_case_study(self, case_study, output_path):
        """
        Export a case study to PowerPoint using the template
        
        Args:
            case_study: CaseStudy object
            output_path: Path to save the generated PowerPoint
        """
        # Load the template
        prs = Presentation(self.template_path)
        
        # Build replacement dictionary
        replacements = self._build_replacements(case_study)
        
        # Process all slides
        for slide in prs.slides:
            # Process shapes in slide
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    self._replace_text_in_shape(shape, replacements)
                
                # Handle tables
                if hasattr(shape, "table"):
                    self._replace_text_in_table(shape.table, replacements)
        
        # Save the presentation
        prs.save(output_path)
        return output_path
    
    def _build_replacements(self, case_study):
        """Build dictionary of placeholder replacements"""
        replacements = {}
        
        for placeholder, field in self.PLACEHOLDERS.items():
            value = getattr(case_study, field, '')
            
            # Format the value
            if value is None:
                value = 'N/A'
            elif field == 'duration_months' and value:
                value = f"{value} months"
            elif field == 'team_size' and value:
                value = f"{value} people"
            elif field == 'project_year' and value:
                value = str(value)
            else:
                value = str(value)
            
            # Add with double curly braces
            replacements[f"{{{{{placeholder}}}}}"] = value
        
        # Add current date
        replacements["{{EXPORT_DATE}}"] = datetime.now().strftime("%B %d, %Y")
        
        return replacements
    
    def _replace_text_in_shape(self, shape, replacements):
        """Replace placeholders in a shape's text"""
        if not shape.has_text_frame:
            return
        
        text_frame = shape.text_frame
        
        # Process each paragraph
        for paragraph in text_frame.paragraphs:
            # Process each run
            for run in paragraph.runs:
                original_text = run.text
                
                # Replace all placeholders
                for placeholder, value in replacements.items():
                    if placeholder in original_text:
                        original_text = original_text.replace(placeholder, value)
                
                run.text = original_text
    
    def _replace_text_in_table(self, table, replacements):
        """Replace placeholders in a table"""
        for row in table.rows:
            for cell in row.cells:
                original_text = cell.text
                
                # Replace all placeholders
                for placeholder, value in replacements.items():
                    if placeholder in original_text:
                        original_text = original_text.replace(placeholder, value)
                
                cell.text = original_text
    
    @staticmethod
    def get_placeholder_guide():
        """Return a guide for available placeholders"""
        guide = """
# PowerPoint Template Placeholder Guide

Use the following placeholders in your PowerPoint template. 
The system will automatically replace them with case study data.

## Available Placeholders:

- {{PROJECT_NAME}} - Name of the project
- {{CLIENT}} - Client name
- {{INDUSTRY}} - Industry sector
- {{YEAR}} - Project year
- {{CHALLENGE}} - Project challenge/problem
- {{SOLUTION}} - Solution provided
- {{OUTCOMES}} - Project outcomes/results
- {{TECHNOLOGIES}} - Technologies used
- {{TEAM_SIZE}} - Team size (e.g., "5 people")
- {{DURATION}} - Project duration (e.g., "6 months")
- {{PROJECT_VALUE}} - Project value/budget
- {{TAGS}} - Project tags
- {{CREATED_BY}} - Created by user
- {{EXPORT_DATE}} - Current date (auto-generated)

## Usage Example:

In your PowerPoint template, add text like:

"Project: {{PROJECT_NAME}}"
"Client: {{CLIENT}}"
"Challenge: {{CHALLENGE}}"

When you export a case study, these will be replaced with actual data.
"""
        return guide
