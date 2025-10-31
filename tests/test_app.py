import io
import os
import shutil
import tempfile
import unittest
from datetime import datetime

from pptx import Presentation

from app import create_app, db


class TestConfig:
    TESTING = True
    BASE_DIR = os.path.abspath(os.path.dirname(__file__) + '/../')
    SQLALCHEMY_DATABASE_URI = 'sqlite:///:memory:'
    SQLALCHEMY_TRACK_MODIFICATIONS = False
    MAX_CONTENT_LENGTH = 16 * 1024 * 1024
    SECRET_KEY = 'test-secret'
    
    @staticmethod
    def init_app(app):
        """Create upload directories used by the application in tests."""
        upload = getattr(TestConfig, 'UPLOAD_FOLDER', os.path.join(TestConfig.BASE_DIR, 'uploads'))
        template = getattr(TestConfig, 'TEMPLATE_FOLDER', os.path.join(upload, 'templates'))
        attachments = getattr(TestConfig, 'ATTACHMENTS_FOLDER', os.path.join(upload, 'attachments'))
        os.makedirs(upload, exist_ok=True)
        os.makedirs(template, exist_ok=True)
        os.makedirs(attachments, exist_ok=True)


class CaseStudyAppTestCase(unittest.TestCase):
    def setUp(self):
        # Temporary upload directories
        self.tempdir = tempfile.mkdtemp()
        TestConfig.UPLOAD_FOLDER = self.tempdir
        TestConfig.TEMPLATE_FOLDER = os.path.join(self.tempdir, 'templates')
        TestConfig.ATTACHMENTS_FOLDER = os.path.join(self.tempdir, 'attachments')
        os.makedirs(TestConfig.TEMPLATE_FOLDER, exist_ok=True)
        os.makedirs(TestConfig.ATTACHMENTS_FOLDER, exist_ok=True)

        # Create the app with test config
        self.app = create_app(TestConfig)
        self.client = self.app.test_client()

        # Create DB
        with self.app.app_context():
            db.create_all()

    def tearDown(self):
        # Remove temp dirs
        shutil.rmtree(self.tempdir, ignore_errors=True)

    def _create_case_payload(self, idx=1):
        return {
            'project_name': f'Project {idx}',
            'client_name': f'Client {idx}',
            'industry': 'Healthcare',
            'project_year': 2024,
            'challenge': 'Legacy systems',
            'solution': 'Cloud migration',
            'outcomes': 'Reduced costs',
            'technologies': 'AWS,Python',
            'team_size': 5,
            'duration_months': 6,
            'tags': 'cloud,migration',
            'project_value': '$100K-$200K',
            'confidential': False,
            'created_by': 'Tester'
        }

    def test_crud_case_study(self):
        # Create
        resp = self.client.post('/api/case-studies', json=self._create_case_payload())
        self.assertEqual(resp.status_code, 201)
        data = resp.get_json()
        cs_id = data['id']

        # Read list
        resp = self.client.get('/api/case-studies')
        self.assertEqual(resp.status_code, 200)
        items = resp.get_json()
        self.assertTrue(any(i['id'] == cs_id for i in items))

        # Read single
        resp = self.client.get(f'/api/case-studies/{cs_id}')
        self.assertEqual(resp.status_code, 200)
        single = resp.get_json()
        self.assertEqual(single['project_name'], 'Project 1')

        # Update
        resp = self.client.put(f'/api/case-studies/{cs_id}', json={'project_name': 'Project 1 Updated'})
        self.assertEqual(resp.status_code, 200)
        updated = resp.get_json()
        self.assertEqual(updated['project_name'], 'Project 1 Updated')

        # Delete
        resp = self.client.delete(f'/api/case-studies/{cs_id}')
        self.assertEqual(resp.status_code, 204)

        # Confirm deletion
        resp = self.client.get(f'/api/case-studies/{cs_id}')
        self.assertEqual(resp.status_code, 404)

    def test_template_upload_and_export(self):
        # Build a minimal pptx in memory
        ppt_path = os.path.join(self.tempdir, 'sample.pptx')
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = '{{PROJECT_NAME}}'
        prs.save(ppt_path)

        # Upload template and set as default
        with open(ppt_path, 'rb') as f:
            data = {
                'name': 'Default Template',
                'description': 'Test template',
                'is_default': 'true'
            }
            file_data = {'file': (f, 'sample.pptx')}
            # Flask test client expects data as dict combining fields and files
            multi = {**data, 'file': (f, 'sample.pptx')}
            resp = self.client.post('/api/templates', data=multi, content_type='multipart/form-data')

        self.assertIn(resp.status_code, (200, 201))
        tpl = resp.get_json()
        self.assertTrue(tpl['is_default'])
        tpl_id = tpl['id']

        # Create a case study to export
        resp = self.client.post('/api/case-studies', json=self._create_case_payload(idx=2))
        self.assertEqual(resp.status_code, 201)
        cs = resp.get_json()
        cs_id = cs['id']

        # Export to pptx (should return a file)
        resp = self.client.get(f'/api/case-studies/{cs_id}/export/pptx')
        # Either 200 with file or 400 if template missing; assert 200 here
        self.assertEqual(resp.status_code, 200)
        # Check mimetype for pptx
        self.assertEqual(resp.mimetype, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
        # Ensure response has content (some bytes)
        self.assertTrue(len(resp.data) > 100)

    def test_stats_and_facets(self):
        # Create multiple case studies
        for i in range(3):
            resp = self.client.post('/api/case-studies', json=self._create_case_payload(idx=i+1))
            self.assertEqual(resp.status_code, 201)

        # Stats
        resp = self.client.get('/api/stats')
        self.assertEqual(resp.status_code, 200)
        data = resp.get_json()
        self.assertIn('total', data)
        self.assertEqual(data['total'], 3)
        self.assertIn('recent', data)

        # Facets (search service)
        resp = self.client.get('/api/facets')
        self.assertEqual(resp.status_code, 200)
        facets = resp.get_json()
        # facets should be a dict-like structure
        self.assertIsInstance(facets, dict)


if __name__ == '__main__':
    unittest.main()
